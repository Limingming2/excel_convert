from urllib.parse import quote

from flask import Flask, render_template, request, redirect, url_for, flash, make_response, send_file
from werkzeug.utils import secure_filename
from markupsafe import escape
import os
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__),'upload/')
ALLOWED_EXTENSIONS = ['xlsx', 'xls']

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


@app.route("/")
def hello_world():
    return "<p>Hello, World</p>"


def show_user_profile():
    # return f'User {escape(username)}'
    return render_template('hello.html')


def allowed_file(filename):
    result = Path(filename).suffix.replace('.', '') in ALLOWED_EXTENSIONS
    print(f"result：{Path(filename).suffix} {result}")
    return result


@app.route('/upload/<filename>', methods=['GET'])
def file_download(filename):
    filename=UPLOAD_FOLDER+filename
    response = make_response(send_file(filename))
    basename = os.path.basename(filename)
    response.headers["Content-Disposition"] = \
        "attachment;" \
        "filename*=UTF-8''{utf_filename}".format(
            utf_filename=quote(basename.encode('utf-8'))
        )
    return response


@app.route('/upload/excel', methods=['GET', 'POST'])
def upload_file():
    print(f'test:{request.method}')
    if request.method == 'POST':
        if 'file' not in request.files:
            print(f'file not in {request.files}')
            flash('No file part')
            return redirect(request.url)
        file = request.files['file']
        print(f'upload filename: {file.filename} type: {Path(file.filename).suffix}')
        if file.filename == '':
            print(f'no selected file')
            flash('No selected file')
            return redirect(request.url)
        if allowed_file(file.filename):

            # filename = secure_filename(file.filename)
            filename = file.filename
            print(f'folder: {UPLOAD_FOLDER}')
            print(f'filename: {filename}')
            excel_file = os.path.join(UPLOAD_FOLDER, filename)
            excel_path = Path(excel_file)
            if excel_path.exists():
                os.remove(excel_file)
            file.save(excel_file)
            ppt_file = create_ppt_with_table(excel_file)
            ppt_file = os.path.basename(ppt_file)
            return (f'<a href=>首页</a>'
                    f'</br>'
                    f'</br>'
                    f'<a href="{ppt_file}">下载ppt</a>')
    return render_template('hello.html')


def create_ppt_with_table(excel_file):
    base_name, ext = os.path.splitext(excel_file)
    ppt_file = f'{base_name}.pptx'
    ppt_path = Path(ppt_file)
    if ppt_path.exists():
        os.remove(ppt_file)
    df = pd.read_excel(excel_file)
    # 创建PPT对象
    prs = Presentation()
    # 按照第一列的序号升序排序
    df_sorted = df.sort_values(by=df.columns[0], ascending=True)
    dic = {}
    for item in df.itertuples():
        key = f'{item[1]}'
        if key not in dic:
            dic[key] = []
        itemlist = dic[key]
        itemlist.append(item)
    # 按页处理数据
    # while i < len(df_sorted):
    for index, item in enumerate(dic):
        print(f'{index}--{item}--{dic[item]}')
        rows_per_slide = len(dic[item])
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
        print(f'创建{index}')

        rows = rows_per_slide  # 当前页显示的行数
        cols = len(df_sorted.columns)  # 列数
        print(f'行{rows}列{cols}')
        # 在幻灯片上创建一个表格
        table = slide.shapes.add_table(rows + 1, cols, Inches(0.4), Inches(1), Inches(14), Inches(5)).table
        table.columns[0].width = Inches(1)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(5)
        table.columns[3].width = Inches(1)
        for col_num, col_name in enumerate(df_sorted.columns):
            table.cell(0, col_num).text = f'{col_name}'  # 表头从第一行开始
        for col_index, column_name in enumerate(df.columns):
            cell = table.cell(0, col_index)
            cell.text = f'{column_name}'
            # 设置表头格式
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    # run.font.bold = True
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(255, 255, 255)  # 白色文本

        # 添加表格

        for i in range(0, len(dic[item])):
            # print(len(dic[item][i]))
            # 创建新的幻灯片
            # 填充表格内容
            for row_num, row in enumerate(dic[item]):
                for col_num, value in enumerate(row):
                    if col_num == 0:
                        continue
                    print(f'第{row_num}行--第{col_num}列--内容：{str(value)}')
                    table.cell(row_num+1, col_num - 1).text = str(value)
        for row_index, row in enumerate(dic[item]):
            for col_index, value in enumerate(row):
                if col_index == 0:
                    continue
                cell = table.cell(row_index + 1, col_index - 1)  # +1 跳过表头
                cell.text = str(value)
                # 设置单元格文本格式
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色文本
    prs.save(ppt_file)
    return ppt_file


if __name__ == "__main__":
    excel_file = "upload/批量箱单.xlsx"  # 输入Excel文件路径
    create_ppt_with_table(excel_file)