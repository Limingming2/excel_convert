import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
# def load_excel_data(excel_file):
#     # 读取Excel文件
#     df = pd.read_excel(excel_file)
#     return df


def create_ppt_with_table(excel_file, ppt_file):
    df = pd.read_excel(excel_file)
    # 创建PPT对象
    prs = Presentation()
    # 按照第一列的序号升序排序
    df_sorted = df.sort_values(by=df.columns[0], ascending=True)
    dic = {}
    for item in df.itertuples():
        key = f'{item[1]}'
        if key not in dic:
            dic[key] = []
        itemlist = dic[key]
        itemlist.append(item)
    # 按页处理数据
    # while i < len(df_sorted):
    for index, item in enumerate(dic):
        print(f'{index}--{item}--{dic[item]}')
        rows_per_slide = len(dic[item])
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
        print(f'创建{index}')

        rows = rows_per_slide  # 当前页显示的行数
        cols = len(df_sorted.columns)  # 列数
        print(f'行{rows}列{cols}')
        # 在幻灯片上创建一个表格
        table = slide.shapes.add_table(rows + 1, cols, Inches(0.4), Inches(1), Inches(14), Inches(5)).table
        table.columns[0].width = Inches(1)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(5)
        table.columns[3].width = Inches(1)
        for col_num, col_name in enumerate(df_sorted.columns):
            table.cell(0, col_num).text = col_name  # 表头从第一行开始
        for col_index, column_name in enumerate(df.columns):
            cell = table.cell(0, col_index)
            cell.text = column_name
            # 设置表头格式
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    # run.font.bold = True
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(255, 255, 255)  # 白色文本

        # 添加表格

        for i in range(0, len(dic[item])):
            # print(len(dic[item][i]))
            # 创建新的幻灯片
            # 填充表格内容
            for row_num, row in enumerate(dic[item]):
                for col_num, value in enumerate(row):
                    if col_num == 0:
                        continue
                    print(f'第{row_num}行--第{col_num}列--内容：{str(value)}')
                    table.cell(row_num+1, col_num - 1).text = str(value)
        for row_index, row in enumerate(dic[item]):
            for col_index, value in enumerate(row):
                if col_index == 0:
                    continue
                cell = table.cell(row_index + 1, col_index - 1)  # +1 跳过表头
                cell.text = str(value)
                # 设置单元格文本格式
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色文本
    prs.save(ppt_file)


if __name__ == "__main__":
    excel_file = "/Users/lmm/Downloads/批量箱单.xlsx"  # 输入Excel文件路径
    ppt_file = "/Users/lmm/Downloads/output_ppt_with_table.pptx"  # 输出PPT文件路径

    # 加载Excel数据
    # df = load_excel_data(excel_file)

    # 创建PPT文件
    create_ppt_with_table(excel_file, ppt_file)

    print("PPT文件已生成，且包含表格！")
